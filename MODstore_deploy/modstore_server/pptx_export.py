"""Generate simple PowerPoint decks from markdown outlines."""

from __future__ import annotations

import io
import re
from typing import List, Tuple


def _parse_markdown_slides(markdown: str, fallback_title: str) -> List[Tuple[str, List[str]]]:
    text = (markdown or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not text:
        return [(fallback_title or "AI 生成 PPT", ["暂无内容"])]

    slides: List[Tuple[str, List[str]]] = []
    current_title = ""
    bullets: List[str] = []

    def flush() -> None:
        nonlocal current_title, bullets
        if current_title or bullets:
            slides.append((current_title or fallback_title or "页面", bullets[:8] or [" "]))
        current_title = ""
        bullets = []

    for raw in text.split("\n"):
        line = raw.strip()
        if not line:
            continue
        h = re.match(r"^#{1,3}\s+(.+)$", line)
        if h:
            flush()
            current_title = h.group(1).strip()[:80]
            continue
        item = re.sub(r"^([-*+]|\d+[.)])\s+", "", line).strip()
        if item:
            bullets.append(item[:180])
    flush()
    return slides[:30] or [(fallback_title or "AI 生成 PPT", ["暂无内容"])]


def build_pptx_from_markdown(markdown: str, *, title: str = "AI 生成 PPT") -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as e:  # pragma: no cover - depends on deployment extras
        raise RuntimeError("服务器未安装 python-pptx，暂不能生成 PPT 文件") from e

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = _parse_markdown_slides(markdown, title)

    # Cover
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    bg = cover.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(10, 12, 20)
    box = cover.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.6), Inches(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title or slides[0][0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 250, 252)
    sub = cover.shapes.add_textbox(Inches(2.2), Inches(4.0), Inches(8.9), Inches(0.6))
    sub.text_frame.text = "由 MODstore AI 创作生成"
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub.text_frame.paragraphs[0].font.size = Pt(16)
    sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(165, 180, 252)

    for idx, (slide_title, bullet_lines) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(15, 23, 42)

        badge = slide.shapes.add_textbox(Inches(0.75), Inches(0.42), Inches(1.0), Inches(0.3))
        badge.text_frame.text = f"{idx:02d}"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(125, 211, 252)

        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.82), Inches(11.8), Inches(0.72))
        title_box.text_frame.text = slide_title
        title_p = title_box.text_frame.paragraphs[0]
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(248, 250, 252)

        body = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.1), Inches(4.8))
        tf = body.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, line in enumerate(bullet_lines[:7]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(226, 232, 240)
            p.space_after = Pt(8)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

